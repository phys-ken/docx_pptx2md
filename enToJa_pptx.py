import requests
from pptx import Presentation
from time import sleep
import os

def translate(str_in, source="en", target="ja"):
    url="https://script.google.com/macros/s/"
    url += "AKfycbzXCTx8lSF5-zn_uvgIbYPKORioVV423NAY995U4OSQ5yf8JjF_8mOVvXvSVUmeRsltPg"
    url += "/exec?text=" + str_in
    url += "&source=" + source
    url += "&target=" + target
    rr = requests.get(url)
    return rr.text

def pptxEn2Ja(filePath,num = 0):
    path_to_presentation = filePath
    prs = Presentation(path_to_presentation)

    print("start___" + filePath)
    for ns, slide in enumerate(prs.slides):
        for nsh, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for np, paragraph in enumerate(shape.text_frame.paragraphs):
                for rs, run in enumerate(paragraph.runs):
                    str_in = run.text
                    str_out = translate(str_in)
                    prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = str_out
                    sleep(1.5)
                    #print(np)
    newFileName = os.path.basename(filePath)
    dirname = "E:/Google(個人)/bad.physics.24のデータ_2018_12_6/0.2021授業用/Pi_ja"
    newFilePath = dirname +  "/"  + str(num).zfill(4) + "_Ja_"+ newFileName
    prs.save(newFilePath)
    os.remove(filePath)
    print(str(num) + "_____end_____" + newFileName)

count = 0
for current, subfolders, subfiles in os.walk("."):
    for fileName in subfiles:
        base, ext = os.path.splitext(fileName)
        if ext == '.pptx' or ext == '.PPTX':
            fullPath = current + "/" + fileName
            print(fullPath)
            count = count + 1
            pptxEn2Ja(fullPath , count)
        else:
            print("無視しました。" + fileName )